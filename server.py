#!/usr/bin/env python3
"""
课件生成 Web 服务
提供网页界面 + PPT 生成 API

依赖：
    pip install flask python-pptx

运行：
    python server.py
    然后访问 http://localhost:5000
"""

import os
import uuid
import json
from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__, template_folder='templates')
UPLOAD_FOLDER = '/Users/qm/.openclaw/workspace/generated'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


def parse_outline(outline_text):
    """解析大纲文本为幻灯片列表"""
    slides = []
    lines = outline_text.strip().split('\n')
    current_section = None
    
    for line in lines:
        trimmed = line.strip()
        if not trimmed:
            continue
        
        if trimmed.startswith('## '):
            current_section = trimmed[3:]
            slides.append({'type': 'section', 'title': current_section})
        elif trimmed.startswith('### '):
            subsection = trimmed[3:]
            slides.append({'type': 'subsection', 'title': subsection, 'content': ''})
        elif trimmed.startswith('# '):
            current_section = trimmed[2:]
            slides.append({'type': 'section', 'title': current_section})
        elif current_section:
            slides.append({'type': 'content', 'title': current_section, 'content': trimmed})
        else:
            slides.append({'type': 'content', 'title': '内容', 'content': trimmed})
    
    return slides


def create_ppt(title, outline_text, output_path):
    """生成 PPT 文件"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle.text = "课件生成"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # 解析大纲并生成内容页
    slides = parse_outline(outline_text)
    
    for slide_data in slides:
        if slide_data['type'] == 'section':
            # 章节页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            title_shape.text_frame.paragraphs[0].font.size = Pt(40)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # 删除内容占位符
            if slide.placeholders[1]:
                sp = slide.placeholders[1]._element
                sp.getparent().remove(sp)
        else:
            # 内容页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            if slide_data.get('content'):
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = slide_data['content']
                tf.paragraphs[0].font.size = Pt(18)
    
    prs.save(output_path)
    return len(prs.slides)


@app.route('/')
def index():
    return render_template('courseware.html')


@app.route('/api/generate', methods=['POST'])
def generate():
    data = request.get_json()
    title = data.get('title', '课件')
    outline = data.get('outline', '')
    
    if not outline:
        return jsonify({'success': False, 'message': '请输入大纲'})
    
    filename = f"{title}_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    
    try:
        slide_count = create_ppt(title, outline, filepath)
        return jsonify({
            'success': True,
            'slides': slide_count,
            'download_url': f'/api/download/{filename}',
            'filename': filename
        })
    except Exception as e:
        return jsonify({'success': False, 'message': str(e)})


@app.route('/api/download/')
def download(filename):
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({'error': '文件不存在'}), 404


if __name__ == '__main__':
    print("🚀 课件生成服务已启动")
    print("   访问 http://localhost:5000")
    app.run(port=5000, debug=True)