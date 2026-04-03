#!/usr/bin/env python3
"""
课件生成 MCP 服务
功能：解析文本内容生成 PowerPoint 课件

依赖：
    pip install python-pptx mcp

用法：
    python courseware_mcp.py
"""

import json
import sys
from typing import Any, Literal
from dataclasses import dataclass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


@dataclass
class Tool:
    name: str
    description: str
    inputSchema: dict


class CoursewareMCPServer:
    """课件生成 MCP 服务"""
    
    def __init__(self):
        self.tools = [
            Tool(
                name="create_powerpoint",
                description="根据教学大纲内容创建 PowerPoint 课件",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "课件标题"},
                        "outline": {"type": "string", "description": "课件大纲（支持多级标题，每行一个标题，##开头为章，###开头为节）"},
                        "output_path": {"type": "string", "description": "输出文件路径"}
                    },
                    "required": ["title", "outline", "output_path"]
                }
            ),
            Tool(
                name="add_slide",
                description="添加一张幻灯片",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "content": {"type": "string", "description": "幻灯片内容（标题|内容格式）"},
                        "layout": {"type": "string", "enum": ["title", "content", "two_content", "blank"], "default": "content"}
                    },
                    "required": ["content"]
                }
            )
        ]
    
    def run(self):
        """MCP 主循环"""
        while True:
            line = sys.stdin.readline()
            if not line:
                break
            
            try:
                request = json.loads(line.strip())
                self.handle_request(request)
            except Exception as e:
                self.send_error(str(e))
    
    def handle_request(self, request: dict):
        """处理 MCP 请求"""
        method = request.get("method")
        
        if method == "tools/list":
            self.send_tools_list()
        elif method == "tools/call":
            self.handle_tool_call(request)
        else:
            self.send_error(f"Unknown method: {method}")
    
    def send_response(self, result: Any):
        """发送响应"""
        print(json.dumps({"result": result}), flush=True)
    
    def send_error(self, error: str):
        """发送错误"""
        print(json.dumps({"error": error}), flush=True)
    
    def send_tools_list(self):
        """发送工具列表"""
        tools = [{"name": t.name, "description": t.description, "inputSchema": t.inputSchema} for t in self.tools]
        self.send_response(tools)
    
    def handle_tool_call(self, request: dict):
        """处理工具调用"""
        name = request.get("params", {}).get("name")
        args = request.get("params", {}).get("arguments", {})
        
        if name == "create_powerpoint":
            result = self.create_powerpoint(**args)
        elif name == "add_slide":
            result = self.add_slide(**args)
        else:
            result = {"error": f"Unknown tool: {name}"}
        
        self.send_response(result)
    
    def create_powerpoint(self, title: str, outline: str, output_path: str) -> dict:
        """创建 PowerPoint 课件"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 封面页
        self.add_title_slide(prs, title)
        
        # 解析大纲并添加内容
        current_section = None
        for line in outline.strip().split('\n'):
            line = line.strip()
            if not line:
                continue
            
            if line.startswith('## '):
                # 新章节
                current_section = line[3:]
                self.add_section_slide(prs, current_section)
            elif line.startswith('### '):
                # 新小节
                section_title = line[3:]
                self.add_content_slide(prs, section_title, "请点击添加内容...")
            elif line.startswith('# '):
                # 新标题
                section_title = line[2:]
                self.add_content_slide(prs, section_title, "请点击添加内容...")
            elif current_section:
                self.add_content_slide(prs, current_section, line)
            else:
                self.add_content_slide(prs, title, line)
        
        # 保存
        prs.save(output_path)
        return {"status": "success", "output_path": output_path}
    
    def add_title_slide(self, prs: Presentation, title: str):
        """添加封面页"""
        slide_layout = prs.slide_layouts[0]  # title slide
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = "课件生成"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    def add_section_slide(self, prs: Presentation, title: str):
        """添加章节页"""
        slide_layout = prs.slide_layouts[1]  # title and content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 删除内容占位符
        if slide.placeholders[1]:
            sp = slide.placeholders[1]._element
            sp.getparent().remove(sp)
    
    def add_content_slide(self, prs: Presentation, title: str, content: str = ""):
        """添加内容页"""
        slide_layout = prs.slide_layouts[1]  # title and content
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # 内容
        if content:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = content
            tf.paragraphs[0].font.size = Pt(18)
    
    def add_slide(self, content: str, layout: str = "content") -> dict:
        """��加单张幻灯片（供后续扩展使用）"""
        return {"status": "pending", "message": "需保存状态后使用"}


def main():
    server = CoursewareMCPServer()
    server.run()


if __name__ == "__main__":
    main()